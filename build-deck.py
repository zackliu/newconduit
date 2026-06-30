"""Generate a 5-slide executive deck for Agent Runtime Sidecar."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Palette
INK = RGBColor(0x10, 0x14, 0x20)
SLATE = RGBColor(0x4A, 0x55, 0x68)
ACCENT = RGBColor(0x2D, 0x6C, 0xDF)
ACCENT2 = RGBColor(0x12, 0xA5, 0x94)
LIGHT = RGBColor(0xF3, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x7A, 0x86, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    r = slide.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
    return r


def band(slide, x, y, w, h, color):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0, space_after=4):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp; p.space_after = Pt(space_after)
        r = p.add_run(); r.text = text; f = r.font
        f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = "Segoe UI"
    return tb


def bullet(slide, x, y, w, h, title, body, tcolor=INK):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = tcolor; r.font.name = "Segoe UI"
    p.space_after = Pt(2)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(12.5); r2.font.color.rgb = SLATE; r2.font.name = "Segoe UI"; p2.line_spacing = 1.05
    return tb


def box(slide, x, y, w, h, fill, border=None):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if border:
        r.line.color.rgb = border; r.line.width = Pt(1.25)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def arrow(slide, x, y, w, color=ACCENT, h=Inches(0.34)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a

# ---------------- Slide 1: Positioning ----------------
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)
band(s, 0, Inches(5.45), SW, Inches(0.12), ACCENT)
txt(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.6),
    [("AGENT RUNTIME SIDECAR", 16, ACCENT2, True)])
txt(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.6), [
    ("An AI Session Management", 50, WHITE, True),
    ("Architecture & Service", 50, WHITE, True),
], sp=1.0)
txt(s, Inches(0.9), Inches(4.1), Inches(11.0), Inches(1.0),
    [("Run stateful, interactive agents as durable online services — where the session, not the machine, is the durable identity.", 20, RGBColor(0xC6,0xD0,0xE0), False)])
txt(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.8), [
    ("Infrastructure for teams putting agents in production: durable sessions · recovery · real-time control · multi-tenant auth & audit", 13, MUTED, False)])

# ---------------- Slide 2: We are a framework, not an app ----------------
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
band(s, 0, 0, SW, Inches(1.25), LIGHT)
txt(s, Inches(0.9), Inches(0.3), Inches(11.5), Inches(0.7), [("Infrastructure, not an end-user application", 27, INK, True)])
txt(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.7), [("No end user opens our product. We are the runtime teams build their agentic applications on — solving the common problems once, instead of in every project.", 14.5, SLATE, False)])
# small aside: not the consumer remote-control product
band(s, Inches(0.9), Inches(2.35), Inches(11.55), Inches(0.7), INK)
txt(s, Inches(1.15), Inches(2.35), Inches(11.0), Inches(0.7), [("Not a consumer tool for operating a single agent — a framework for the team that builds the application.", 13.5, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
# the problems you hit
txt(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.4), [("PROBLEMS EVERY AGENTIC-APPLICATION TEAM FACES", 13, ACCENT, True)])
probs = [
    ("Session continuity", "A long-running agent is one continuous job, not a stateless request."),
    ("Reconnect", "Clients disconnect and return to the same session, on any machine."),
    ("Failure recovery", "When compute fails, the session is recovered rather than lost."),
    ("Multi-tenant & audit", "Tenant isolation and a full record of access on every path."),
]
cw, gx = Inches(2.78), Inches(0.2); x0, y0 = Inches(0.9), Inches(3.85); ch = Inches(2.9)
for i,(t,b) in enumerate(probs):
    cx = x0+(cw+gx)*i
    box(s, cx, y0, cw, ch, LIGHT); band(s, cx, y0, cw, Inches(0.1), ACCENT if i%2==0 else ACCENT2)
    bullet(s, cx+Inches(0.28), y0+Inches(0.35), cw-Inches(0.5), ch-Inches(0.5), t, b)

# ---------------- Slide 3: Scenario — building an ops agent system ----------------
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
band(s, 0, 0, SW, Inches(1.25), LIGHT)
txt(s, Inches(0.9), Inches(0.3), Inches(11.5), Inches(0.7), [("Scenario: building an operations / incident agent system", 25, INK, True)])
txt(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.6), [("Each step demands a capability. Different tasks require different isolated environments \u2014 each defined as an AgentSpec.", 14.5, SLATE, False)])
# flow stages
sy, bh, bw = Inches(2.25), Inches(2.35), Inches(2.78)
xs = [Inches(0.9), Inches(4.0), Inches(7.1), Inches(10.2)]
stages = [
    ("Incident raised", "An alert opens a case; the application starts a session for the incident.", ACCENT, "trigger"),
    ("Troubleshooting agent", "Runs diagnostic skills and tools / MCP in an isolated environment.", ACCENT2, "AgentSpec A"),
    ("Code-fix agent", "Clones the repository to build and test \u2014 a separate environment, a separate AgentSpec.", ACCENT, "AgentSpec B"),
    ("Verify & resolve", "Resumes later, continues the prior work, and returns a result.", ACCENT2, "continue"),
]
for i,(t,b,col,tag) in enumerate(stages):
    box(s, xs[i], sy, bw, bh, LIGHT); band(s, xs[i], sy, bw, Inches(0.1), col)
    txt(s, xs[i]+Inches(0.25), sy+Inches(0.3), bw-Inches(0.45), Inches(0.4), [(tag.upper(), 11, col, True)])
    txt(s, xs[i]+Inches(0.25), sy+Inches(0.7), bw-Inches(0.45), Inches(0.6), [(t, 16, INK, True)])
    txt(s, xs[i]+Inches(0.25), sy+Inches(1.3), bw-Inches(0.45), Inches(1.0), [(b, 12, SLATE, False)], sp=1.05)
    if i < 3:
        arrow(s, xs[i]+bw+Inches(0.02), sy+bh/2-Inches(0.17), Inches(0.28), MUTED)
# runtime band underneath
ry = Inches(5.1)
band(s, Inches(0.9), ry, Inches(12.08), Inches(1.55), INK)
txt(s, Inches(1.2), ry+Inches(0.22), Inches(11.6), Inches(0.45), [("ONE SESSION RUNTIME CARRIES THE WORK END TO END", 12, ACCENT2, True)])
txt(s, Inches(1.2), ry+Inches(0.65), Inches(11.6), Inches(0.8), [
    ("Durable session per incident  \u00b7  isolation by AgentSpec  \u00b7  pause while waiting  \u00b7  reconnect  \u00b7  recover and continue on a new worker", 14, WHITE, False)])

# ---------------- Slide 4: Architecture (layered) ----------------
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
band(s, 0, 0, SW, Inches(1.25), LIGHT)
txt(s, Inches(0.9), Inches(0.3), Inches(11.5), Inches(0.7), [("Architecture: one source of truth, pluggable layers", 25, INK, True)])
txt(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5), [("Central owns the session. Controllers translate protocols; managers own workflows; adapters bind technology. Each adapter is replaceable.", 13.5, SLATE, False)])
def vline(x, y1, y2):
    c = s.shapes.add_connector(2, x, y1, x, y2); c.line.color.rgb = MUTED; c.line.width = Pt(1.5)
# Consumers
for i,(t,sub) in enumerate([("Clients","browser · desktop · CLI"),("Application backends","create & manage sessions"),("Workers","sidecar + agent process")]):
    cx = Inches(0.9)+Inches(4.18)*i
    box(s, cx, Inches(1.95), Inches(3.95), Inches(0.78), RGBColor(0xE8,0xEE,0xFA))
    txt(s, cx, Inches(2.0), Inches(3.95), Inches(0.7), [(t,13.5,INK,True),(sub,10.5,SLATE,False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
# SDK boundary
band(s, Inches(0.9), Inches(2.95), Inches(11.55), Inches(0.45), ACCENT)
txt(s, Inches(0.9), Inches(2.95), Inches(11.55), Inches(0.45), [("SDK / API  ·  HTTPS / WebSocket",12,WHITE,True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
vline(Inches(6.67), Inches(2.73), Inches(2.95))
# Central control plane
box(s, Inches(0.9), Inches(3.6), Inches(8.05), Inches(2.05), INK)
txt(s, Inches(1.15), Inches(3.7), Inches(7.6), Inches(0.4), [("CENTRAL SESSION SERVICE — control plane",12.5,ACCENT2,True)])
for j,c in enumerate(["Client ingress","Worker ingress","Agent events"]):
    bx=Inches(1.15)+Inches(2.5)*j; box(s,bx,Inches(4.15),Inches(2.35),Inches(0.5),RGBColor(0x22,0x2B,0x3D)); txt(s,bx,Inches(4.15),Inches(2.35),Inches(0.5),[(c,11,WHITE,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for j,c in enumerate(["Session","Registry","Routing","Leases","Snapshot","Audit"]):
    bx=Inches(1.15)+Inches(1.25)*j; box(s,bx,Inches(4.8),Inches(1.15),Inches(0.7),RGBColor(0x22,0x2B,0x3D)); txt(s,bx,Inches(4.8),Inches(1.15),Inches(0.7),[(c,10.5,RGBColor(0xC6,0xD0,0xE0),True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# Storage
box(s, Inches(9.15), Inches(3.6), Inches(3.3), Inches(2.05), RGBColor(0xE8,0xEE,0xFA), border=ACCENT)
txt(s, Inches(9.35), Inches(3.8), Inches(2.95), Inches(1.7), [("Persistent storage",13,INK,True),("session catalog",11,SLATE,False),("event log · snapshots",11,SLATE,False),("artifacts · audit",11,SLATE,False)], anchor=MSO_ANCHOR.MIDDLE, space_after=4)
# Pluggable adapters
txt(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.35), [("PLUGGABLE ADAPTERS",12,ACCENT,True)])
for i,(t,v) in enumerate([("Transport","Web PubSub"),("Hosting","Docker"),("Storage","local files"),("Agent","Copilot SDK")]):
    cx=Inches(0.9)+Inches(2.95)*i; box(s,cx,Inches(6.25),Inches(2.8),Inches(0.85),LIGHT,border=ACCENT2)
    txt(s,cx,Inches(6.25),Inches(2.8),Inches(0.85),[(t,12.5,INK,True),(v,11,MUTED,False)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=1)


# ---------------- Slide 5: Why it matters ----------------
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)
band(s, 0, Inches(1.4), SW, Inches(0.1), ACCENT)
txt(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.8), [("Why this is worth building", 34, WHITE, True)])
points = [
    ("Every team rebuilds this layer", "Session routing, worker registry, event log, snapshots, and auth are reimplemented in every project."),
    ("We provide the missing runtime", "Not the model and not the framework \u2014 the durable session layer between them."),
    ("Differentiated and defensible", "The value is in the operational runtime, not the model; existing agents are adopted through a sidecar."),
    ("Self-host first, managed later", "One model from local proof of concept to production cluster to managed service."),
]
y = Inches(1.9)
for t,b in points:
    band(s, Inches(0.9), y+Inches(0.1), Inches(0.12), Inches(0.9), ACCENT2)
    txt(s, Inches(1.2), y, Inches(11.0), Inches(1.0), [(t,19,WHITE,True),(b,13.5,RGBColor(0xC6,0xD0,0xE0),False)], space_after=2)
    y += Inches(1.18)

import os, time
out = r"c:\Users\chenyl\newconduit\Agent-Runtime-Sidecar.pptx"
try:
    prs.save(out)
except PermissionError:
    out = rf"c:\Users\chenyl\newconduit\Agent-Runtime-Sidecar-v{int(time.time())}.pptx"
    prs.save(out)
print("saved", out)
